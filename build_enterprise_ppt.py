import sys
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Modern Dark Enterprise Theme
    BG_DARK = RGBColor(15, 23, 42)       # Slate 900
    BG_CARD = RGBColor(30, 41, 59)       # Slate 800
    TEXT_LIGHT = RGBColor(248, 250, 252) # Slate 50
    TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
    ACCENT_RED = RGBColor(225, 29, 72)   # Rose 600
    ACCENT_BLUE = RGBColor(14, 165, 233) # Sky 500
    ACCENT_TEAL = RGBColor(20, 184, 166) # Teal 500
    ACCENT_AMBER = RGBColor(245, 158, 11)# Amber 500
    BORDER_COLOR = RGBColor(51, 65, 85)  # Slate 700

    blank_layout = prs.slide_layouts[6]

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_header(slide, title_text, category_text="SANJEEVANI (AERO) ENTERPRISE ARCHITECTURE"):
        # Category Banner
        tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = tb_cat.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_BLUE

        # Title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_LIGHT

    def add_card(slide, left, top, width, height, title, body_bullets, accent_color=ACCENT_BLUE):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = BORDER_COLOR
        shape.line.width = Pt(1)

        # Header Accent Line inside card
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.2), Inches(top + 0.2), Inches(width - 0.4), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = accent_color
        line.line.fill.background()

        # Text inside card
        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.3), Inches(width - 0.4), Inches(height - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_LIGHT
        p_t.space_after = Pt(8)

        for b in body_bullets:
            p = tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(4)

    # ── SLIDE 1: Title Slide ───────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)

    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "SANJEEVANI (AERO)"
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_RED
    p0.space_after = Pt(10)

    p1 = tf.add_paragraph()
    p1.text = "Enterprise Event-Driven Emergency Response Platform & System Architecture"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    p1.space_after = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "Stateless AI Advisory  •  Emergency Orchestrator  •  Deterministic Resource Allocation  •  Hospital Gateway Adapters"
    p2.font.size = Pt(13)
    p2.font.color.rgb = ACCENT_BLUE

    # ── SLIDE 2: Executive Summary & CTO Redesign ───────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_header(slide2, "Executive Redesign: Event-Driven & Service-Oriented Architecture")

    add_card(slide2, 0.8, 1.6, 5.7, 5.2, "1. Emergency Orchestrator (The Brain)", [
        "Replaces rigid AI-centric flows with a central Event Orchestrator.",
        "Manages emergency lifecycle: CREATED ➔ TRIAGED ➔ DISPATCHED ➔ RESOLVED.",
        "Coordinates stateless AI services, matching engine, and notification gateway.",
        "Nothing communicates directly; all events route through Orchestrator."
    ], ACCENT_RED)

    add_card(slide2, 6.8, 1.6, 5.7, 5.2, "2. Decoupled Enterprise Core Services", [
        "Stateless AI Service: Returns advisory JSON (type, severity, confidence) without dispatching.",
        "Deterministic Allocation Engine: Multi-factor mathematical scoring (Beds, ICU, Doctor, Distance, ETA, Blood).",
        "Hospital Gateway Adapter: Pluggable REST, HIS, FHIR & Dashboard adapters.",
        "Async Event Bus & Notification Service: Real-time WebSockets, Push FCM, SMS & WhatsApp."
    ], ACCENT_TEAL)

    # ── SLIDE 3: System Topology & Component UML Diagram ────────────────
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_header(slide3, "System Architecture Topology & Modular Subsystems")

    add_card(slide3, 0.8, 1.6, 3.7, 5.2, "Client & Intake Layer", [
        "Citizen PWA & Native App: Audio recording & real-time voice streaming.",
        "Staff Platform: Unified app for Doctors, Drivers & ASHA Workers.",
        "Hospital Command Portal: Real-time bed & emergency queue management.",
        "FastAPI REST Gateway: Central API router with JWT security."
    ], ACCENT_BLUE)

    add_card(slide3, 4.8, 1.6, 3.7, 5.2, "Core Orchestrator & AI Engine", [
        "Emergency Orchestrator: Central event controller & state machine.",
        "Whisper v3 STT: Auto-detects 12+ Indian languages & dialects.",
        "Groq Llama 3.3 70B: Sub-second intent & clinical triage engine.",
        "Gemini 2.0 & Rule Fallback: Zero-downtime multi-tier redundancy."
    ], ACCENT_RED)

    add_card(slide3, 8.8, 1.6, 3.7, 5.2, "Gateway & Infrastructure Layer", [
        "Hospital Adapter Resolver: Universal API client & Dashboard gateway.",
        "Deterministic Matching Engine: Mathematical scoring matrix.",
        "Protocol Engine: Evidence-based WHO / AHA first-aid retrieval.",
        "PostgreSQL + Supabase: Unified relational persistence & cloud storage."
    ], ACCENT_AMBER)

    # ── SLIDE 4: Deterministic Resource Allocation Engine ───────────────
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_header(slide4, "Deterministic Resource Allocation Engine (Judgable Scoring Matrix)")

    add_card(slide4, 0.8, 1.6, 5.7, 5.2, "Multi-Factor Mathematical Inputs", [
        "Patient Clinical Severity: RED (1.0), AMBER (0.7), GREEN (0.3).",
        "Geospatial Proximity & Live Traffic ETA: Real-time GPS distance calculation.",
        "Real-Time Hospital Capacity: Total Beds, ICU Beds, Ventilators & Blood Bank status.",
        "On-Duty Medical Roster: Available Specialists, Surgeons & Emergency Physicians.",
        "ASHA/ANM Rural Proximity: Local community responder ETA comparison."
    ], ACCENT_AMBER)

    add_card(slide4, 6.8, 1.6, 5.7, 5.2, "Deterministic Output Scoring", [
        "Hospital Score = f(Specialty, Beds, Distance, ICU Availability).",
        "Ambulance Score = f(Vehicle Type, Live GPS ETA, Driver Shift).",
        "ASHA Score = f(Local Distance, Skill Set, Training Category).",
        "100% Explainable & Auditable: Zero black-box AI dispatch risks; fully compliant with health regulations."
    ], ACCENT_TEAL)

    # ── SLIDE 5: Hospital Gateway & Adapter Pattern ─────────────────────
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_header(slide5, "Hospital Gateway: Pluggable Adapter Pattern")

    add_card(slide5, 0.8, 1.6, 5.7, 5.2, "Tier A — Corporate & Large Hospitals", [
        "REST / Webhook API Adapter: Direct integration with Apollo, Yashoda, Care, AIG HIS.",
        "Zero Workflow Disruption: Inbound webhook triggers internal hospital dispatch.",
        "HL7 / FHIR Adapter: Standardized healthcare interoperability data exchange.",
        "Automated Roster & Bed Sync: Automated real-time capacity telemetry."
    ], ACCENT_BLUE)

    add_card(slide5, 6.8, 1.6, 5.7, 5.2, "Tier B — Small Hospitals & CHCs", [
        "Free Sanjeevani Dashboard Adapter: Additive capability for facilities without IT infra.",
        "Accept / Decline / Escalate: One-click emergency case acceptance workflow.",
        "Manual Capacity Toggles: Simple UI controls for bed and doctor availability.",
        "Zero Integration Friction: Instant onboarding for rural and small town clinics."
    ], ACCENT_RED)

    # ── SLIDE 6: Split AI Pipeline & Evidence-Based Protocol Engine ─────
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)
    add_header(slide6, "Decoupled AI Pipeline & Safety-Certified Protocol Engine")

    add_card(slide6, 0.8, 1.6, 5.7, 5.2, "Split AI Advisory Pipeline", [
        "Speech Service: Groq Whisper Large v3 (transcribes 12+ Indian scripts).",
        "Translation Service: Converts regional speech into clinical English.",
        "Intent Service: Groq Llama 3.3 70B classifies sub-intent (Cardiac, Trauma, etc.).",
        "Clinical Triage Service: Severity rating & chief complaint extraction.",
        "Stateless Advisory: AI provides recommendations; humans & deterministic engines execute."
    ], ACCENT_TEAL)

    add_card(slide6, 6.8, 1.6, 5.7, 5.2, "Safety-Certified Protocol Engine", [
        "Evidence-Based First Aid: Governed strictly by WHO, AHA, and Red Cross guidelines.",
        "Zero Hallucination Guarantee: AI never invents medical or first-aid instructions.",
        "Multilingual Script Rendering: Renders steps in Telugu, Hindi, Tamil, Bengali, etc.",
        "Native TTS Audio Guidance: Spoken step-by-step guidance directly on caller screen."
    ], ACCENT_AMBER)

    # ── SLIDE 7: Event Bus & Notification Service ───────────────────────
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7)
    add_header(slide7, "Async Event Bus & Centralized Notification Service")

    add_card(slide7, 0.8, 1.6, 5.7, 5.2, "Event Bus Architecture", [
        "Asynchronous Event Emitters: EmergencyCreated, EmergencyTriaged, ResponderAssigned, EmergencyAccepted.",
        "Decoupled Subscribers: Hospital Gateway, Staff Apps, Command Center & Audit Service subscribe independently.",
        "Zero Cascading Failures: High availability even if individual channels experience latency."
    ], ACCENT_RED)

    add_card(slide7, 6.8, 1.6, 5.7, 5.2, "Central Notification Router", [
        "Multi-Channel Dispatch: WebSocket, FCM Push, SMS, WhatsApp & Voice Call.",
        "Intelligent Channel Selection: Sends push to app; falls back to SMS/WhatsApp for low network.",
        "Real-Time Family Broadcast: Automatically alerts emergency contacts with live tracking link."
    ], ACCENT_BLUE)

    # ── SLIDE 8: Emergency Lifecycle State Machine & Incident Audit ─────
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8)
    add_header(slide8, "Emergency State Machine & Full Incident Audit Trail")

    add_card(slide8, 0.8, 1.6, 5.7, 5.2, "Strict State Machine Transitions", [
        "CREATED ➔ ANALYZING: Raw report received via voice/text.",
        "TRIAGED ➔ MATCHING: Severity & sub-intent established by AI.",
        "WAITING_ACCEPTANCE ➔ ACCEPTED: Hospital or responder confirms case.",
        "DISPATCHED ➔ EN_ROUTE: Ambulance moving to patient GPS.",
        "ARRIVED ➔ RESOLVED: Patient admitted to ER; case completed."
    ], ACCENT_AMBER)

    add_card(slide8, 6.8, 1.6, 5.7, 5.2, "Immutable Incident Audit Log", [
        "Millisecond Timestamps: Every state transition and API call logged.",
        "Structured Hospital Handover: Generates complete digital handover summary for ER doctor.",
        "Legal & Regulatory Compliance: Full auditability for healthcare governance.",
        "Performance Analytics: Heatmaps, response time metrics, and hospital acceptance rates."
    ], ACCENT_TEAL)

    # ── SLIDE 9: Multi-Role Platform Suite & Command Center ─────────────
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9)
    add_header(slide9, "Multi-Role Application Suite & Command Center")

    add_card(slide9, 0.8, 1.6, 3.7, 5.2, "Citizen SOS App", [
        "PWA & Native Mobile Web.",
        "One-tap voice & text intake.",
        "Whisper v3 12+ Indian language STT.",
        "Live GPS radar map animation.",
        "Native speech first-aid playback."
    ], ACCENT_RED)

    add_card(slide9, 4.8, 1.6, 3.7, 5.2, "Staff & Responder App", [
        "Unified app for Doctors, Drivers & ASHA.",
        "Role-based dynamic UI layout.",
        "One-click accept / decline dispatch.",
        "Turn-by-turn navigation & Siren alert.",
        "Live bed & shift toggle."
    ], ACCENT_BLUE)

    add_card(slide9, 8.8, 1.6, 3.7, 5.2, "Hospital & Command Portal", [
        "Tier B Hospital Management System.",
        "District Command Center GIS Map.",
        "Real-time ICU/Bed & Doctor capacity.",
        "Super Admin verification dashboard.",
        "Analytics & Regional Heatmaps."
    ], ACCENT_TEAL)

    # ── SLIDE 10: Tech Stack & Production Deployment ─────────────────────
    slide10 = prs.slides.add_slide(blank_layout)
    set_background(slide10)
    add_header(slide10, "Technology Stack & Production Architecture")

    add_card(slide10, 0.8, 1.6, 5.7, 5.2, "Core Application Stack", [
        "Backend: FastAPI (Python 3.12/3.14) + SQLModel ORM + Pydantic v2.",
        "Database: PostgreSQL on Supabase + PostGIS geospatial extensions.",
        "Web Frontend: React + Vite + Vanilla CSS Glassmorphism Design System.",
        "Mobile App: Flutter (Cross-platform iOS / Android)."
    ], ACCENT_BLUE)

    add_card(slide10, 6.8, 1.6, 5.7, 5.2, "AI & Cloud Integration Stack", [
        "Speech-to-Text: Groq Whisper Large v3 (sub-second multilingual STT).",
        "LLM Intelligence: Groq Llama 3.3 70B (0.05s response time).",
        "Backup AI: Google Gemini 2.0 Flash Lite & Rule Engine.",
        "Cloud Storage: Supabase Storage Bucket (`hospital-documents`)."
    ], ACCENT_RED)

    # ── SLIDE 11: Summary & Architecture Superiority ────────────────────
    slide11 = prs.slides.add_slide(blank_layout)
    set_background(slide11)

    tb = slide11.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "THE ARCHITECTURE IS THE ARGUMENT."
    p0.font.size = Pt(38)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_RED
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "Sanjeevani (AERO) scales without asking any responder to change first."
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    p1.space_after = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "✓ Event-Driven Emergency Orchestrator\n✓ Stateless AI Advisory Engine\n✓ Deterministic Resource Allocation Matrix\n✓ Pluggable Hospital Gateway Adapters\n✓ 100% Explainable, Safe & Auditable Health Platform"
    p2.font.size = Pt(14)
    p2.font.color.rgb = ACCENT_TEAL

    prs.save("Sanjeevani_Enterprise_Architecture.pptx")
    prs.save("Sanjeevani_Round2_Architecture.pptx")
    print("Successfully generated Sanjeevani_Enterprise_Architecture.pptx and updated Sanjeevani_Round2_Architecture.pptx!")

if __name__ == "__main__":
    create_presentation()
